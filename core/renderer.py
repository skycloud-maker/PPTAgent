"""Render a SlideSchema into a company-internal PPTX deck.

LG전자 내부 보고 장표 스타일 기준:
- 표지: 좌측 L자 레드 장식 + 중앙 제목/부제 + 하단 부서|날짜 + 우하단 로고 + CONFIDENTIAL
- 본문: 좌상단 제목 + 검은 실선 + 상단 중앙 'LGE Internal Use Only'
- 풋노트: 초록색(#006600) 부연설명
- 강조: LG 레드(#A50034)
- 기본: 검정
- 한글 폰트: LG스마트체 / 영어·기호: Arial Narrow
"""

from __future__ import annotations

import io
import logging
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from core.schema import SlideContent, SlideSchema, SlideType

logger = logging.getLogger(__name__)

# ── 환경변수 ──────────────────────────────────────────────────
COMPANY_NAME       = os.getenv("PPTAGENT_COMPANY_NAME", "LG Electronics")
CONFIDENTIAL_LABEL = os.getenv("PPTAGENT_CONFIDENTIAL_LABEL", "LGE Internal Use Only")

# ── LG 브랜드 색상 ────────────────────────────────────────────
COLOR_RED      = RGBColor(0xA5, 0x00, 0x34)   # LG 레드 (Primary / 강조)
COLOR_BLACK    = RGBColor(0x00, 0x00, 0x00)   # 기본 텍스트
COLOR_GREEN    = RGBColor(0x00, 0x66, 0x00)   # 풋노트 (초록)
COLOR_GRAY     = RGBColor(0x55, 0x55, 0x55)   # 보조 텍스트
COLOR_GRAY_LT  = RGBColor(0xF2, 0xF2, 0xF2)  # 배경 박스
COLOR_BORDER   = RGBColor(0xD9, 0xD9, 0xD9)  # 구분선
COLOR_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# ── 폰트 ─────────────────────────────────────────────────────
FONT_KO  = os.getenv("PPTAGENT_FONT_KO", "LG스마트체")   # 한글
FONT_EN  = os.getenv("PPTAGENT_FONT_EN", "Arial Narrow") # 영어·기호·숫자

# ── 슬라이드 규격 16:9 ────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.6)
MARGIN_R = Inches(0.6)

# ── 로고 ─────────────────────────────────────────────────────
_ASSETS_DIR     = os.path.normpath(os.path.join(os.path.dirname(__file__), "..", "assets"))
COMPANY_LOGO    = os.path.join(_ASSETS_DIR, "company_logo.png")
LOGO_PATH       = os.getenv("PPTAGENT_LOGO_PATH", COMPANY_LOGO)


# ─────────────────────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────────────────────

def render_pptx(schema: SlideSchema) -> bytes:
    """SlideSchema → .pptx bytes (디스크 저장 없음 — SECURITY.md)"""
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for sd in schema.slides:
        slide = prs.slides.add_slide(blank)
        logger.info("Rendering | index=%s | type=%s", sd.index, sd.type.value)

        if sd.type == SlideType.TITLE:
            _render_title(slide, sd.content)
        elif sd.type == SlideType.SECTION:
            _render_section(slide, sd.content)
        elif sd.type == SlideType.BULLET:
            _render_bullet(slide, sd.content)
        elif sd.type == SlideType.CHART:
            _render_chart_placeholder(slide, sd.content)
        elif sd.type == SlideType.TABLE:
            _render_table_placeholder(slide, sd.content)
        elif sd.type == SlideType.TWO_COLUMN:
            _render_two_column(slide, sd.content)
        elif sd.type == SlideType.IMAGE:
            _render_image(slide, sd.content)
        else:
            _render_default(slide, sd.content)

        # 표지 제외 공통 헤더/푸터
        if sd.type != SlideType.TITLE:
            _add_header(slide)
            _add_footer(slide, sd.index)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────
# 공통 고정 요소
# ─────────────────────────────────────────────────────────────

def _add_header(slide) -> None:
    """상단 공통: 중앙 'LGE Internal Use Only' + 우상단 로고(또는 회사명)"""
    # 중앙 상단 기밀 문구
    _tb(slide,
        l=Inches(4.5), t=Inches(0.08),
        w=Inches(4.3), h=Inches(0.28),
        text=CONFIDENTIAL_LABEL,
        font=FONT_EN, size=9,
        color=COLOR_GRAY, align=PP_ALIGN.CENTER)

    # 우상단 로고 또는 회사명
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(
            LOGO_PATH,
            left=Inches(12.1), top=Inches(0.08),
            width=Inches(0.95), height=Inches(0.38))
    else:
        _tb(slide,
            l=Inches(11.5), t=Inches(0.08),
            w=Inches(1.6), h=Inches(0.28),
            text=COMPANY_NAME,
            font=FONT_EN, size=9, bold=True,
            color=COLOR_RED, align=PP_ALIGN.RIGHT)


def _add_title_bar(slide, title_text: str) -> None:
    """좌상단 제목 + 바로 아래 검은 실선 (본문 슬라이드 공통)"""
    _tb(slide,
        l=MARGIN_L, t=Inches(0.38),
        w=Inches(11.5), h=Inches(0.52),
        text=title_text,
        font=FONT_KO, size=20, bold=True,
        color=COLOR_BLACK)

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        MARGIN_L, Inches(0.94),
        Inches(12.1), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_BLACK
    line.line.fill.background()


def _add_footer(slide, page_num: int) -> None:
    """하단 중앙 페이지 번호"""
    _tb(slide,
        l=Inches(6.2), t=Inches(7.1),
        w=Inches(1.0), h=Inches(0.28),
        text=str(page_num),
        font=FONT_EN, size=9,
        color=COLOR_GRAY, align=PP_ALIGN.CENTER)


def _add_footnote(slide, text: str, top: float = 6.7) -> None:
    """초록색 풋노트 텍스트 추가"""
    _tb(slide,
        l=MARGIN_L, t=Inches(top),
        w=Inches(12.1), h=Inches(0.35),
        text=text,
        font=FONT_KO, size=10,
        color=COLOR_GREEN)


# ─────────────────────────────────────────────────────────────
# 슬라이드 타입별 렌더러
# ─────────────────────────────────────────────────────────────

def _render_title(slide, content: SlideContent) -> None:
    """
    표지 슬라이드 — 실제 LG 장표 스타일
    - 상단 중앙: LGE Internal Use Only
    - 좌측: L자 레드 장식 (세로 바 + 가로 바)
    - 중앙: 메인 제목(큰) + 부제(작음)
    - 하단: 부서명 | 날짜
    - 우하단: LG 로고
    - 최하단 중앙: | CONFIDENTIAL |
    """
    # ── 상단 기밀 문구
    _tb(slide,
        l=Inches(4.5), t=Inches(0.1),
        w=Inches(4.3), h=Inches(0.28),
        text=CONFIDENTIAL_LABEL,
        font=FONT_EN, size=9,
        color=COLOR_GRAY, align=PP_ALIGN.CENTER)

    # ── 좌측 L자 레드 장식
    # 세로 바 (좌측 상단 ~ 중간)
    v_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.0), Inches(0.0),
        Inches(0.18), Inches(4.2))
    v_bar.fill.solid()
    v_bar.fill.fore_color.rgb = COLOR_RED
    v_bar.line.fill.background()

    # 가로 바 (L자 하단)
    h_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.0), Inches(4.02),
        Inches(1.5), Inches(0.18))
    h_bar.fill.solid()
    h_bar.fill.fore_color.rgb = COLOR_RED
    h_bar.line.fill.background()

    # ── 메인 제목
    _tb(slide,
        l=Inches(1.8), t=Inches(2.2),
        w=Inches(9.5), h=Inches(1.4),
        text=content.title or "제목을 입력해주세요",
        font=FONT_KO, size=36, bold=True,
        color=COLOR_BLACK)

    # ── 부제목 (안건 / 부제)
    if content.subtitle:
        _tb(slide,
            l=Inches(1.8), t=Inches(3.7),
            w=Inches(9.5), h=Inches(0.6),
            text=content.subtitle,
            font=FONT_KO, size=16,
            color=COLOR_GRAY)

    # ── 하단 부서명 | 날짜
    presenter_text = content.presenter or ""
    if presenter_text:
        _tb(slide,
            l=Inches(1.8), t=Inches(5.6),
            w=Inches(6.0), h=Inches(0.4),
            text=presenter_text,
            font=FONT_KO, size=13,
            color=COLOR_GRAY)

    # ── 우하단 LG 로고
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(
            LOGO_PATH,
            left=Inches(11.3), top=Inches(5.3),
            width=Inches(1.6), height=Inches(0.8))

    # ── 최하단 CONFIDENTIAL
    _tb(slide,
        l=Inches(4.5), t=Inches(7.05),
        w=Inches(4.3), h=Inches(0.28),
        text="| CONFIDENTIAL |",
        font=FONT_EN, size=9,
        color=COLOR_RED, align=PP_ALIGN.CENTER)


def _render_section(slide, content: SlideContent) -> None:
    """섹션 구분 슬라이드"""
    # 좌측 레드 컬럼
    col = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0, 0, Inches(5.0), SLIDE_H)
    col.fill.solid()
    col.fill.fore_color.rgb = COLOR_RED
    col.line.fill.background()

    _tb(slide,
        l=Inches(0.4), t=Inches(2.8),
        w=Inches(4.2), h=Inches(1.8),
        text=content.heading or content.title or "",
        font=FONT_KO, size=26, bold=True,
        color=COLOR_WHITE, align=PP_ALIGN.CENTER)


def _render_bullet(slide, content: SlideContent) -> None:
    """불릿 슬라이드 — 제목 + 실선 + 본문"""
    heading = content.heading or content.title or ""
    points  = content.points or []

    _add_title_bar(slide, heading)

    for i, point in enumerate(points[:5]):
        _tb(slide,
            l=Inches(0.85), t=Inches(1.15 + i * 1.0),
            w=Inches(11.6), h=Inches(0.85),
            text=f"• {point}",
            font=FONT_KO, size=14,
            color=COLOR_BLACK)

    # 풋노트
    if content.caption:
        _add_footnote(slide, content.caption)


def _render_chart_placeholder(slide, content: SlideContent) -> None:
    """차트 슬라이드 (Phase 3에서 실제 차트 구현)"""
    heading = content.heading or content.title or ""
    _add_title_bar(slide, heading)

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.1),
        Inches(11.7), Inches(5.6))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_GRAY_LT
    box.line.color.rgb = COLOR_BORDER

    _tb(slide,
        l=Inches(0.8), t=Inches(3.5),
        w=Inches(11.7), h=Inches(0.6),
        text="[ 차트 영역 ]",
        font=FONT_EN, size=13,
        color=COLOR_BORDER, align=PP_ALIGN.CENTER)

    if content.caption:
        _add_footnote(slide, content.caption)


def _render_table_placeholder(slide, content: SlideContent) -> None:
    """표 슬라이드 (Phase 3에서 실제 표 구현)"""
    heading = content.heading or content.title or ""
    _add_title_bar(slide, heading)

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.1),
        Inches(11.7), Inches(5.6))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_GRAY_LT
    box.line.color.rgb = COLOR_BORDER

    _tb(slide,
        l=Inches(0.8), t=Inches(3.5),
        w=Inches(11.7), h=Inches(0.6),
        text="[ 표 영역 ]",
        font=FONT_EN, size=13,
        color=COLOR_BORDER, align=PP_ALIGN.CENTER)

    if content.caption:
        _add_footnote(slide, content.caption)


def _render_two_column(slide, content: SlideContent) -> None:
    """2단 비교 슬라이드"""
    heading = content.heading or content.title or ""
    _add_title_bar(slide, heading)

    # 중앙 구분선
    div = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(6.55), Inches(1.1),
        Pt(1.2), Inches(5.8))
    div.fill.solid()
    div.fill.fore_color.rgb = COLOR_BORDER
    div.line.fill.background()

    # 좌측 헤더
    left_title  = content.left_title  or "항목 1"
    right_title = content.right_title or "항목 2"

    _tb(slide,
        l=Inches(0.8), t=Inches(1.1),
        w=Inches(5.5), h=Inches(0.45),
        text=left_title,
        font=FONT_KO, size=14, bold=True,
        color=COLOR_RED, align=PP_ALIGN.CENTER)

    _tb(slide,
        l=Inches(7.0), t=Inches(1.1),
        w=Inches(5.5), h=Inches(0.45),
        text=right_title,
        font=FONT_KO, size=14, bold=True,
        color=COLOR_RED, align=PP_ALIGN.CENTER)

    # 좌측 항목
    for i, point in enumerate((content.left_points or [])[:5]):
        _tb(slide,
            l=Inches(0.9), t=Inches(1.7 + i * 0.85),
            w=Inches(5.4), h=Inches(0.7),
            text=f"• {point}",
            font=FONT_KO, size=13,
            color=COLOR_BLACK)

    # 우측 항목
    for i, point in enumerate((content.right_points or [])[:5]):
        _tb(slide,
            l=Inches(7.1), t=Inches(1.7 + i * 0.85),
            w=Inches(5.4), h=Inches(0.7),
            text=f"• {point}",
            font=FONT_KO, size=13,
            color=COLOR_BLACK)

    if content.caption:
        _add_footnote(slide, content.caption)


def _render_image(slide, content: SlideContent) -> None:
    """이미지 슬라이드"""
    heading = content.heading or content.title or ""
    _add_title_bar(slide, heading)

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.1),
        Inches(11.7), Inches(5.6))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_GRAY_LT
    box.line.color.rgb = COLOR_BORDER

    if content.image_path and os.path.exists(content.image_path):
        slide.shapes.add_picture(
            content.image_path,
            left=Inches(1.0), top=Inches(1.3),
            width=Inches(11.3), height=Inches(5.2))
    else:
        _tb(slide,
            l=Inches(0.8), t=Inches(3.5),
            w=Inches(11.7), h=Inches(0.6),
            text="[ 이미지 영역 ]",
            font=FONT_EN, size=13,
            color=COLOR_BORDER, align=PP_ALIGN.CENTER)

    if content.caption:
        _add_footnote(slide, content.caption)


def _render_default(slide, content: SlideContent) -> None:
    """기본 슬라이드"""
    heading = content.heading or content.title or ""
    if heading:
        _add_title_bar(slide, heading)


# ─────────────────────────────────────────────────────────────
# 헬퍼
# ─────────────────────────────────────────────────────────────

def _tb(
    slide,
    l: Emu, t: Emu, w: Emu, h: Emu,
    text: str,
    font: str = FONT_KO,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """텍스트 박스 생성 헬퍼"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color or COLOR_BLACK
